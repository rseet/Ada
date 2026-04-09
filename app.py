import os
import json
import streamlit as st
from dotenv import load_dotenv
from openai import AzureOpenAI
from io import BytesIO
from copy import deepcopy
from pptx import Presentation
from pathlib import Path
from pypdf import PdfReader

# ---------- PPTX TEMPLATE ----------
BASE_DIR = Path(__file__).parent
TEMPLATE_PATH = BASE_DIR / "KNB Research Slides Template 2025.pptx"

if not TEMPLATE_PATH.exists():
    st.error(f"Template not found at: {TEMPLATE_PATH}")
    st.stop()

# ---------- PAGE CONFIG (must be first Streamlit command) ----------
st.set_page_config(page_title="Research First-Cut Workbench", layout="wide")
st.title("Research First-Cut Workbench")
st.caption("From brief to defensible storyboard - before PowerPoint.")

# ---------- ENV / AZURE CLIENT ----------
if Path(".env").exists():
    load_dotenv()

AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION", "2024-02-15-preview")
AZURE_OPENAI_DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT")  # must match Azure Deployment Name

# Optional debug (remove once stable)
st.write("AZURE_OPENAI_API_KEY loaded:", bool(AZURE_OPENAI_API_KEY))
st.write("AZURE_OPENAI_ENDPOINT:", AZURE_OPENAI_ENDPOINT)
st.write("AZURE_OPENAI_DEPLOYMENT:", AZURE_OPENAI_DEPLOYMENT)

client = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    api_key=AZURE_OPENAI_API_KEY,
    api_version=AZURE_OPENAI_API_VERSION
)

# ---------- SESSION STATE ----------
if "brief" not in st.session_state:
    st.session_state.brief = None
if "synthesis" not in st.session_state:
    st.session_state.synthesis = None
if "storyboard" not in st.session_state:
    st.session_state.storyboard = None

# ---------- HELPER: CALL AI ----------
def call_ai(system_prompt, user_prompt, temperature=0.4):
    try:
        response = client.chat.completions.create(
            model=AZURE_OPENAI_DEPLOYMENT,  # deployment name (e.g. gpt-4.1-mini)
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            temperature=temperature,
            response_format={"type": "json_object"},
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        st.error(f"AI call failed: {e}")
        return None


st.write("OPENAI_API_KEY loaded?", bool(os.getenv("OPENAI_API_KEY")))


# ---HELPER FUNCTIONS FOR SLIDEMAKING --------

def remove_slide(prs, slide):
    """Remove a slide from a python-pptx Presentation."""
    slide_id_list = prs.slides._sldIdLst  # pylint: disable=protected-access
    slides = list(prs.slides)
    idx = slides.index(slide)
    slide_id = slide_id_list[idx].rId
    slide_id_list.remove(slide_id_list[idx])
    prs.part.drop_rel(slide_id)

def duplicate_slide(prs, source_slide):
    """
    Duplicate slide within the same presentation by copying XML shapes.
    Good for text-based slides (no charts). Images usually OK; complex media may vary.
    """
    blank_layout = prs.slide_layouts[6]  # usually blank
    new_slide = prs.slides.add_slide(blank_layout)

    # Copy shapes
    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # Copy relationships (helps with images/icons sometimes)
    for rel in source_slide.part.rels.values():
        # skip notesSlide relationships
        if "notesSlide" not in rel.reltype:
            try:
                new_slide.part.rels.add_relationship(rel.reltype, rel._target, rel.rId)
            except Exception:
                pass

    return new_slide

def iter_text_shapes(slide):
    for shape in slide.shapes:
        if shape.has_text_frame:
            yield shape

def replace_text_contains(slide, needle, replacement):
    """
    Replace entire text in any shape whose text contains 'needle'.
    Useful because template has marker strings like 'Take Away Box' / 'Main Title Area' / 'Contents Area'. [1](https://khazanahcommy-my.sharepoint.com/personal/ruixi_seet_khazanah_com_my/_layouts/15/Doc.aspx?sourcedoc=%7B1ED756FF-422B-4F38-A7F0-E77396F20B2E%7D&file=KNB%20Research%20Slides%20Template%202025.pptx&action=edit&mobileredirect=true)[2](https://khazanahcommy-my.sharepoint.com/personal/ruixi_seet_khazanah_com_my/_layouts/15/Doc.aspx?sourcedoc=%7BDA802A2D-B95C-40ED-95B6-B40FB35F97FE%7D&file=KNB%20Research%20Slides%20Template%202025.pptx&action=edit&mobileredirect=true&DefaultItemOpen=1)
    """
    for shape in iter_text_shapes(slide):
        text = shape.text_frame.text or ""
        if needle in text:
            # Replace all text in the shape (simple & robust)
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = replacement

def set_title(slide, title_text):
    for shape in iter_text_shapes(slide):
        if "Main title" in (shape.text_frame.text or ""):
            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = title_text
            # Remove the built-in title placeholder ("Click to add title") if present
            try:
                title_ph = slide.shapes.title
                if title_ph is not None:
                    title_ph.element.getparent().remove(title_ph.element)
            except Exception:
                pass
            return
    try:
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide.shapes.title.text_frame.clear()
            slide.shapes.title.text_frame.paragraphs[0].text = title_text
    except Exception:
        pass

def add_bullets_to_contents(slide, bullets):
    # 1) try marker
    for shape in iter_text_shapes(slide):
        if "Contents Area" in (shape.text_frame.text or ""):
            shape.text_frame.clear()
            for i, b in enumerate(bullets):
                p = shape.text_frame.paragraphs[0] if i == 0 else shape.text_frame.add_paragraph()
                p.text = b
                p.level = 0
            return

    # 2) fallback: biggest text box by area
    candidates = [s for s in iter_text_shapes(slide)]
    if not candidates:
        return
    biggest = max(candidates, key=lambda s: (s.width * s.height))
    biggest.text_frame.clear()
    for i, b in enumerate(bullets):
        p = biggest.text_frame.paragraphs[0] if i == 0 else biggest.text_frame.add_paragraph()
        p.text = b
        p.level = 0

def add_speaker_notes(slide, note_lines):
    """Add speaker notes (e.g., evidence/source tags)."""
    notes_slide = slide.notes_slide  # creates if missing [6](https://python-pptx-fix.readthedocs.io/en/stable/dev/analysis/sld-notes-slide.html)
    tf = getattr(notes_slide, "notes_text_frame", None)
    if tf is None:
        return  # fail silently; avoids hard crash on edge decks [7](https://github.com/scanny/python-pptx/issues/983)

    tf.clear()
    for i, line in enumerate(note_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line

        
def batch_rewrite_titles_greenboxes(slides, tone, word_limit=12):
    """
    slides: list of dicts containing slideNumber, draftTitle, greenBox
    returns: list of dicts {slideNumber, finalTitle, finalGreenBox}
    """
    system_prompt = f"""You rewrite research slide titles and green boxes.

STRICT RULES:
- Titles: max {word_limit} words
- Green box: max 2–3 short lines (hard limit)
- Do NOT introduce any new facts or claims
- Maintain original meaning exactly
- Be assertive and insight-led
- Tone: suitable for {tone}

Return valid JSON:
{{
  "slides": [
    {{"slideNumber": number, "finalTitle": "string", "finalGreenBox": "string"}}
  ]
}}"""

    user_prompt = "Rewrite the following slides:\n" + json.dumps({"slides": slides}, indent=2)
    result = call_ai(system_prompt, user_prompt, temperature=0.3)
    if not result:
        return None
    return result.get("slides", [])


# ============================================================
# STEP 1: RESEARCH BRIEF
# ============================================================
st.header("Step 1: Research Brief")
st.markdown("*Complete this before using any AI tools.*")

with st.form("brief_form"):
    col1, col2 = st.columns(2)

    with col1:
        deck_objective = st.text_area(
            "Deck Objective",
            placeholder="e.g. Evaluate demand outlook for Saudi residential real estate over next 3 years"
        )
        audience = st.selectbox(
            "Target Audience",
            ["Investment Committee", "Board of Directors", "Portfolio Manager", "Internal Team"]
        )
        core_questions = st.text_area(
            "Core Question(s)",
            placeholder="e.g. Is demand accelerating or decelerating? What are the key drivers?"
        )

    with col2:
        scope = st.text_area(
            "Scope Boundaries",
            placeholder="e.g. Saudi Arabia only, 2022-2027, exclude commercial"
        )
        hypothesis = st.text_area(
            "Initial Hypothesis / Angle (optional)",
            placeholder="e.g. We believe demand is structurally supported by Vision 2030 demographics"
        )
        uploaded_files = st.file_uploader(
            "Upload source materials (optional)",
            accept_multiple_files=True,
            type=["txt", "md","pdf"]
        )

    submitted = st.form_submit_button("Lock Brief")

    if submitted:
        if not deck_objective or not core_questions:
            st.error("Deck Objective and Core Questions are required.")
        else:
            source_texts = []
            uploaded_files = uploaded_files or []
    
            for f in uploaded_files:
                suffix = Path(f.name).suffix.lower()
    
                try:
                    if suffix in [".txt", ".md"]:
                        # Streamlit UploadedFile is bytes
                        text = f.read().decode("utf-8", errors="replace")
    
                    elif suffix == ".pdf":
                        # Ensure pointer at start
                        f.seek(0)
                        reader = PdfReader(f)
                        pages = []
                        for i, page in enumerate(reader.pages):
                            page_text = page.extract_text() or ""
                            pages.append(f"\n--- Page {i+1} ---\n{page_text}")
                        text = "\n".join(pages)
    
                    else:
                        text = "(unsupported file type)"
    
                    source_texts.append(f"[{f.name}]\n{text}")
    
                except Exception as e:
                    source_texts.append(f"[{f.name}] (could not read: {e})")
    
            st.session_state.brief = {
                "deckObjective": deck_objective,
                "audience": audience,
                "coreQuestions": core_questions,
                "scope": scope,
                "hypothesis": hypothesis,
                "sourceTexts": source_texts
            }
            st.success("Brief locked. You can now run AI tools below.")

# Show locked brief
if st.session_state.brief:
    with st.expander("View Locked Brief"):
        st.json(st.session_state.brief)

# ============================================================
# STEP 2: AI TOOLS
# ============================================================
if st.session_state.brief:
    st.header("Step 2: Choose AI Action")
    st.markdown("*Select one tool at a time. You control what runs.*")

    tab1, tab2, tab3, tab4 = st.tabs([
        "Research Synthesiser",
        "Storyboard Generator",
        "Title / Green Box Rewriter",
        "Deck Generator"
    ])

    brief = st.session_state.brief

    # ----------------------------------------------------------
    # TAB 1: RESEARCH SYNTHESISER
    # ----------------------------------------------------------
    with tab1:
        st.subheader("Research Synthesiser")
        st.markdown("Turns your brief and sources into structured, source-tagged research notes.")

        if st.button("Run Synthesiser", key="run_synth"):
            with st.spinner("Synthesising research..."):

                system_prompt = """You are a senior research analyst assistant.
Your job is to synthesise source materials against a research brief.

STRICT RULES:
- Every factual claim MUST have a source tag
- Flag any numbers without clear sources
- Never provide conclusions or recommendations
- Be explicit about confidence and gaps

Return valid JSON in this exact format:
{
  "themes": [
    {
      "theme": "string",
      "evidence": [
        {"text": "string", "sourceTag": "Broker|Consultant|Official|News|Other", "sourceName": "string"}
      ],
      "conflictingViewpoints": [{"view": "string", "source": "string"}],
      "confidence": "High|Medium|Low",
      "dataGaps": ["string"]
    }
  ],
  "overallGaps": ["string"],
  "suggestedFollowUps": ["string"]
}"""

                source_block = "\n\n".join(brief["sourceTexts"]) if brief["sourceTexts"] else "No source materials provided."

                user_prompt = f"""RESEARCH BRIEF:
Objective: {brief['deckObjective']}
Audience: {brief['audience']}
Core Questions: {brief['coreQuestions']}
Scope: {brief['scope']}
Hypothesis: {brief['hypothesis']}

SOURCE MATERIALS:
{source_block}

Please synthesise the above into structured research notes."""

                result = call_ai(system_prompt, user_prompt)

                if result:
                    st.session_state.synthesis = result
                    st.success("Synthesis complete.")

        if st.session_state.synthesis:
            synth = st.session_state.synthesis

            for i, theme in enumerate(synth.get("themes", [])):
                with st.expander(f"Theme {i+1}: {theme['theme']} - Confidence: {theme['confidence']}"):
                    st.markdown("**Evidence:**")
                    for ev in theme.get("evidence", []):
                        st.markdown(f"- {ev['text']} `[{ev['sourceTag']}: {ev.get('sourceName', 'unknown')}]`")

                    if theme.get("conflictingViewpoints"):
                        st.markdown("**Conflicting Viewpoints:**")
                        for cv in theme["conflictingViewpoints"]:
                            st.markdown(f"- {cv['view']} `[{cv['source']}]`")

                    if theme.get("dataGaps"):
                        st.markdown("**Data Gaps:**")
                        for gap in theme["dataGaps"]:
                            st.markdown(f"- {gap}")

            if synth.get("overallGaps"):
                st.markdown("### Overall Data Gaps")
                for g in synth["overallGaps"]:
                    st.markdown(f"- {g}")

            if synth.get("suggestedFollowUps"):
                st.markdown("### Suggested Follow-Ups")
                for f_item in synth["suggestedFollowUps"]:
                    st.markdown(f"- {f_item}")

    # ----------------------------------------------------------
    # TAB 2: STORYBOARD GENERATOR
    # ----------------------------------------------------------
    with tab2:
        st.subheader("Storyboard Generator")
        st.markdown("Generates 2-3 alternative slide structures. **You choose the story.**")

        use_synthesis = st.checkbox(
            "Include synthesis results as input",
            value=st.session_state.synthesis is not None,
            key="use_synth"
        )

        if st.button("Generate Storyboards", key="run_story"):
            with st.spinner("Building storyboard options..."):

                system_prompt = """You are a senior research deck architect.
Your job is to propose 2-3 alternative storyboard structures for a research deck.

STRICT RULES:
- Do NOT pick a best option. Present alternatives equally.
- Each slide must have a clear objective, assertive title, and evidence.
- Flag confidence and gaps per slide.
- Do NOT include formatting, layout, or design.
- Suggest chart TYPES only (e.g. bar, line, waterfall), never actual charts.

Return valid JSON:
{
  "storyflows": [
    {
      "name": "string",
      "description": "string",
      "slides": [
        {
          "slideNumber": 1,
          "slideObjective": "string",
          "draftTitle": "string",
          "greenBox": "string",
          "supportingEvidence": ["string"],
          "suggestedChartType": "string",
          "confidence": "High|Medium|Low",
          "knownGaps": ["string"]
        }
      ]
    }
  ]
}"""

                synth_block = ""
                if use_synthesis and st.session_state.synthesis:
                    synth_block = f"\n\nSYNTHESISED RESEARCH:\n{json.dumps(st.session_state.synthesis, indent=2)}"

                user_prompt = f"""RESEARCH BRIEF:
Objective: {brief['deckObjective']}
Audience: {brief['audience']}
Core Questions: {brief['coreQuestions']}
Scope: {brief['scope']}
Hypothesis: {brief['hypothesis']}
{synth_block}

Generate 2-3 alternative storyboard structures."""

                result = call_ai(system_prompt, user_prompt, temperature=0.6)

                if result:
                    st.session_state.storyboard = result
                    st.success("Storyboards generated.")

        if st.session_state.storyboard:
            flows = st.session_state.storyboard.get("storyflows", [])

            for flow in flows:
                st.markdown(f"### {flow['name']}")
                st.markdown(f"*{flow['description']}*")

                for slide in flow.get("slides", []):
                    with st.expander(
                        f"Slide {slide['slideNumber']}: {slide['draftTitle']} "
                        f"[{slide['confidence']}]"
                    ):
                        st.markdown(f"**Objective:** {slide['slideObjective']}")
                        st.markdown(f"**Title:** {slide['draftTitle']}")
                        st.markdown(f"**Green Box:** {slide['greenBox']}")
                        st.markdown("**Evidence:**")
                        for ev in slide.get("supportingEvidence", []):
                            st.markdown(f"- {ev}")
                        if slide.get("suggestedChartType"):
                            st.markdown(f"**Chart type:** {slide['suggestedChartType']}")
                        if slide.get("knownGaps"):
                            st.markdown("**Gaps:**")
                            for gap in slide["knownGaps"]:
                                st.markdown(f"- {gap}")

                st.divider()

    # ----------------------------------------------------------
    # TAB 3: TITLE / GREEN BOX REWRITER
    # ----------------------------------------------------------
    with tab3:
        st.subheader("Title / Green Box Rewriter")
        st.markdown("Paste any title or green box. Get 3-5 clearer alternatives.")

        input_text = st.text_area(
            "Your title or green box",
            placeholder="e.g. The market is expected to grow due to various factors"
        )
        tone = st.selectbox(
            "Tone",
            ["Investment Committee", "Board of Directors", "Portfolio Manager"],
            key="rewrite_tone"
        )
        word_limit = st.slider("Max words", 5, 25, 12)

        if st.button("Rewrite", key="run_rewrite"):
            if not input_text:
                st.warning("Please enter text to rewrite.")
            else:
                with st.spinner("Rewriting..."):

                    system_prompt = f"""You rewrite research slide titles and green boxes.

STRICT RULES:
- Maximum {word_limit} words per option
- Do NOT introduce any new facts or claims
- Maintain the original meaning exactly
- Be assertive and insight-led
- Tone: suitable for {tone}

Return valid JSON:
{{
  "options": [
    {{"text": "string", "wordCount": number, "clarityNote": "string"}}
  ]
}}"""

                    user_prompt = f'Rewrite this: "{input_text}"'

                    result = call_ai(system_prompt, user_prompt)

                    if result:
                        st.markdown("**Options:**")
                        for i, opt in enumerate(result.get("options", [])):
                            st.markdown(
                                f"{i+1}. **{opt['text']}** "
                                f"({opt['wordCount']} words) - _{opt.get('clarityNote', '')}_"
                            )
    # ----------------------------------------------------------
    # TAB 4: DECK GENERATOR
    # ----------------------------------------------------------
    with tab4:
        st.subheader("Deck Generator")
        st.markdown("Generates a PPTX deck using the selected storyboard + rewritten titles/green boxes.")
    
        if not st.session_state.storyboard:
            st.info("Generate storyboards first (Tab 2).")
        else:
            flows = st.session_state.storyboard.get("storyflows", [])
            flow_names = [f["name"] for f in flows] if flows else []
            selected_name = st.radio("Choose a storyboard flow", flow_names) if flow_names else None
    
            tone = st.selectbox(
                "Tone for rewrite (same as your Tab 3)",
                ["Investment Committee", "Board of Directors", "Portfolio Manager"],
                key="deck_tone"
            )
            deck_word_limit = st.slider("Title max words", 5, 25, 12, key="deck_word_limit")
            include_notes = st.checkbox("Put supporting evidence into speaker notes", value=True)
            include_chart_placeholders = st.checkbox("Add chart-type placeholder text", value=True)
    
            # ---- LOCAL TEMPLATE PATH (NO UPLOADER) ----
            st.caption(f"Using template: {TEMPLATE_PATH}")
    
            if st.button("Generate Deck", key="generate_deck_btn"):
                if not selected_name:
                    st.error("No storyboard flow available.")
                else:
                    try:
                        selected_flow = next((f for f in flows if f["name"] == selected_name), None)
                        slides_in_flow = selected_flow.get("slides", []) if selected_flow else []
    
                        # batch rewrite titles/green boxes
                        rewrite_input = [
                            {"slideNumber": s["slideNumber"], "draftTitle": s["draftTitle"], "greenBox": s["greenBox"]}
                            for s in slides_in_flow
                        ]
                        rewritten = batch_rewrite_titles_greenboxes(
                            rewrite_input, tone=tone, word_limit=deck_word_limit
                        )
    
                        rewrite_map = {x["slideNumber"]: x for x in (rewritten or [])}
                        for s in slides_in_flow:
                            r = rewrite_map.get(s["slideNumber"])
                            s["finalTitle"] = r["finalTitle"] if r else s["draftTitle"]
                            s["finalGreenBox"] = r["finalGreenBox"] if r else s["greenBox"]
    
                        # Build PPTX from LOCAL template
                        prs = Presentation(TEMPLATE_PATH)
    
                        base_cover = prs.slides[0]
                        base_content = prs.slides[2]
                        base_divider = prs.slides[4] if len(prs.slides) > 4 else prs.slides[1]  # fixed ">"
    
                        created_slides = []
    
                        # Cover slide
                        cover = duplicate_slide(prs, base_cover)
                        set_title(cover, brief["deckObjective"][:80])
                        replace_text_contains(cover, "SubTitle", brief.get("coreQuestions", "")[:120])
                        created_slides.append(cover)
    
                        # Content slides
                        for s in slides_in_flow:
                            new_slide = duplicate_slide(prs, base_content)
    
                            set_title(new_slide, s.get("finalTitle", s["draftTitle"]))
                            replace_text_contains(new_slide, "Take Away Box", s.get("finalGreenBox", s["greenBox"]))
    
                            bullets = []
                            bullets.extend(s.get("supportingEvidence", []))
    
                            if include_chart_placeholders and s.get("suggestedChartType"):
                                bullets.append(f"[Chart placeholder: {s['suggestedChartType']}]")
    
                            for g in s.get("knownGaps", []):
                                bullets.append(f"[Gap] {g}")
    
                            if bullets:
                                add_bullets_to_contents(new_slide, bullets)
    
                            if include_notes:
                                note_lines = [
                                    f"Slide {s['slideNumber']} | Confidence: {s.get('confidence','')}",
                                    "Evidence:"
                                ] + [f"- {e}" for e in s.get("supportingEvidence", [])]
    
                                if s.get("knownGaps"):
                                    note_lines += ["", "Known gaps:"] + [f"- {g}" for g in s["knownGaps"]]
    
                                add_speaker_notes(new_slide, note_lines)
    
                            created_slides.append(new_slide)
    
                        # Remove original template slides
                        originals = list(prs.slides)
                        for sl in originals:
                            if sl not in created_slides:
                                try:
                                    remove_slide(prs, sl)
                                except Exception:
                                    pass
    
                        out = BytesIO()
                        prs.save(out)
                        out.seek(0)
    
                        st.success("Deck generated.")
                        st.download_button(
                            label="Download PPTX",
                            data=out,
                            file_name=f"{selected_name.replace(' ', '_')}_deck.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )
                    except Exception as e:
                        st.error(f"Deck generation failed: {e}")
        
    
# ============================================================
# FOOTER
# ============================================================
st.divider()
st.caption(
    "Research First-Cut Workbench - MVP v0.1 | "
    "AI assists structure only. Analyst owns the story.")


